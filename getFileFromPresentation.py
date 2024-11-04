'''Задача 2024.10.25.03
Поменяйте код так, чтобы он выводил текст каждого шэйпа на экран.'''
from pptx import Presentation # подключаем класс Перс... из пптикс библиотеки

# получаем ссылку на файл с презентацией
prs = Presentation('src\Zimnyaya_igra_1.pptx')
# бежим по слайдам
for slide in prs.slides:
    # бежим по шейпам на каждом слайде
    for shape in slide.shapes:
        # если шейп на слайде пустой, то идем к следующему шейпу
        if not shape.has_text_frame:
            continue
        # выводим в консоль содержимое фрейма
        print(shape.text_frame.text)
