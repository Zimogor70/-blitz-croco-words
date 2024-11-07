from pptx import Presentation
from check_spell import check
import os
import datetime


def only_one_word(line: str) -> bool:
    return ' ' not in line and line.isalpha()


def send_to_check(st: set) -> list:
    word_list = list()
    date_list = list()
    str_words: str = ''
    # разбиваем по : и отправляем в списки
    for wd in st:
        sptext = wd.split(':')
        word_list.append(sptext[0])
        date_list.append(sptext[1])

    str_words = ' '.join(word_list)  # concat all to string separ=' '
    ret_list = check(str_words).split(' ')  # str to list

    for i in range(len(ret_list)):
        ret_list[i] += ':' + date_list[i]
    return ret_list


def save(st: set):
    with open('words.txt', 'w', encoding="utf8") as file:
        file.write('')
        file.write('\n'.join(st))


def pars(ls_files) -> set:
    words_set = set()
    for name in ls_files:
        prs = Presentation('src\\'.__add__(name))
        m_timestamp = os.path.getmtime('src\\'.__add__(name))
        m_datestamp = datetime.datetime.fromtimestamp(m_timestamp)
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if only_one_word(shape.text.strip()):
                    words_set.add(shape.text.strip() + ':' + str(m_datestamp.date()))
    return words_set
