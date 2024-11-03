'''
Найдите способ получить содержимое презентации внутри zip файла.
'''
import zipfile
from pptx import Presentation

def open_zip(name: str):
    list_of_pres: list = []
    a: str = ""
    with zipfile.ZipFile(name, 'r') as archive:
        for f in archive.namelist():
            list_of_pres.extend(parsing(archive.open(f)))
    return list_of_pres


def parsing(name):
    words = list()
    prs = Presentation(name)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if ' ' in shape.text or ':' in shape.text or 'БЛИЦ-КРОКОДИЛ' in shape.text:
                continue
            words.append(shape.text)
    return words