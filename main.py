'''Задача 2024.10.30.04
Теперь попробуйте достать все презентации прямо из архива.
То есть ваша программа должна работать, когда вы поулчаете на вход zip архив,
остаёте из него все pptx файлы и парсите каждый, сохраняете из него все слова,
 и после парсинга всех презентаций сохраняете все слова в words.txt'''
import zipfile
from pptx import Presentation


def extr_zip(name: str) -> list:
    list_of_pres: list = []
    a: str = ""
    with zipfile.ZipFile(name, 'r') as archive:
        for f in archive.namelist():
            list_of_pres.append(f)
        archive.extractall('src\\')
    return list_of_pres


# def parsing(name):
#     words = list()
#     prs = Presentation(name)
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if not shape.has_text_frame:
#                 continue
#             if ' ' in shape.text or ':' in shape.text or 'БЛИЦ-КРОКОДИЛ' in shape.text:
#                 continue
#             words.append(shape.text)
#     return words


def only_one_word(line: str) -> bool:
    return ' ' not in line and line.isalpha()


def parsing(ls_files):
    words_list = list('')
    with open('words.txt', 'w', encoding="utf8") as file:
        file.write('')
        for name in ls_files:
            prs = Presentation('src\\'.__add__(name))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    if only_one_word(shape.text.strip()):
                        words_list.append(shape.text.strip() + '\n')
            file.writelines(words_list)


def main():
    parsing(extr_zip('src\croco-blitz-source.zip'))


if __name__ == "__main__":
    main()
